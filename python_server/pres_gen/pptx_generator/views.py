from django.http.response import HttpResponsePermanentRedirect
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt
from io import BytesIO
from pptx import Presentation
import json
import re

@csrf_exempt
def index(request):

    text = json.loads(request.body)["test"]
    app = PresentationConverter(text)
    app.process_text()
    prs = app.create_presentation()

    target_stream = BytesIO()
    prs.save(target_stream)
    response = HttpResponse(content_type='application/vnd.ms-powerpoint')
    response['Content-Disposition'] = 'attachment; filename="sample.pptx"'
    source_stream = BytesIO()
    prs.save(source_stream)
    ppt = source_stream.getvalue()
    source_stream.close()
    response.write(ppt)
    return response


class PresentationConverter: 

    def __init__(self, text):
        self.prs = Presentation()
        self.prs_width = self.prs.slide_width
        self.prs_height = self.prs.slide_height
        self.text = text

    '''
    Process the text and create blocks of text corresponding to individual slides
    '''
    def process_text(self):
        text_array = re.split('\\n{2,}', self.text) # split the text file into slides
        self.slide_chunks = [] # create the "chunks" -> each chunk corresponds to a different slide
        if (len(text_array) == 1 and text_array[0] == ''): # empty presentation
            return
        # title slide --
        title_slide_text = text_array[0].split('\n')
        self.slide_chunks.append({
            'title': title_slide_text[0],
            'subtitle': title_slide_text[1]
        })
        # other slides --
        for i in range(1, len(text_array)):
            slide_text = text_array[i].split('\n', 1)
            slide_chunk = {}
            if (len(slide_text) == 0): # empty slide
                continue
            slide_chunk["title"] = slide_text[0]
            if (len(slide_text) == 1): # slide with only title
                continue
            # slide_chunk["subtitle"] = slide_text[1]
            # slide_chunk["bullets"] = []
            # for i in range(2, len(slide_text)):
            #     slide_chunk["bullets"].append(slide_text[i])
            if (slide_text[1][0:7] == 'image: '): # if there is an image
                image_line_and_bullets_split = slide_text[1].split('\n', 1)
                slide_chunk["image"] = image_line_and_bullets_split[0][7:]
                if (len(image_line_and_bullets_split) > 1): # determine if there are any bullet points (if this is an image-only slide)
                    slide_chunk["bullets"] = image_line_and_bullets_split[1]
            else:
                slide_chunk["bullets"] = slide_text[1]
            self.slide_chunks.append(slide_chunk)
            
    '''
    Create the presentation using the text chunks created in process_text
    '''
    def create_presentation(self):

        # Create title slide --
        title_lyt = self.prs.slide_layouts[0] # get layout
        title_slide = self.prs.slides.add_slide(title_lyt) # create title slide with layout
        title_title = title_slide.shapes.title # get title from title slide
        title_subtitle = title_slide.placeholders[1] # get subtitle from title slide

        if (len(self.slide_chunks) == 0): # empty presentation
            self.prs.save('slide1.pptx')
            return 

        title_chunk = self.slide_chunks[0] 

        title_title.text = title_chunk["title"] # set title
        title_subtitle.text = title_chunk["subtitle"] # set subtitle

        # Create other slides --
        slide_lyt = self.prs.slide_layouts[1] # regular slide layout
        image_slide_lyt = self.prs.slide_layouts[5] # image slide layout

        for i in range(1, len(self.slide_chunks)):
            slide_chunk = self.slide_chunks[i]
            if ("image" in slide_chunk):
                if ("bullets" in slide_chunk): # image slide with text
                    slide_slide = self.prs.slides.add_slide(slide_lyt) # add slide
                    slide_title = slide_slide.shapes.title # get title from slide
                    slide_bullets = slide_slide.placeholders[1] # get bullet points area from slide
                    slide_title.text = slide_chunk["title"] # set title
                    left = self.prs_width * 0.80
                    top = self.prs_height * 0.01
                    width=self.prs_width * 0.16
                    try:
                        slide_slide.shapes.add_picture(slide_chunk["image"], left, top, width=width) # add the image to top right corner
                    except FileNotFoundError as err: # if image path is invalid, tell user
                        print('Image was not found: ' + slide_chunk["image"])
                    slide_bullets.text = slide_chunk["bullets"] # set bullet points
                else: # image only slide
                    image_slide_slide = self.prs.slides.add_slide(image_slide_lyt) # add slide
                    image_slide_title = image_slide_slide.shapes.title # get title
                    image_slide_title.text = slide_chunk["title"] # set title
                    top = self.prs_height*0.20
                    width = self.prs_width*0.6
                    left = self.prs_width*0.2
                    try:
                        image_slide_slide.shapes.add_picture(slide_chunk["image"], left, top, width=width) # add the image to middle of slide
                    except FileNotFoundError as err: # if image path is invalid, tell user
                        print('Image was not found: ' + slide_chunk["image"])
            else: 
                slide_slide = self.prs.slides.add_slide(slide_lyt) # add slide
                slide_title = slide_slide.shapes.title # get title from slide
                slide_bullets = slide_slide.placeholders[1] # get bullet points area from slide
                slide_title.text = slide_chunk["title"] # set ttile
                slide_bullets.text = slide_chunk["bullets"] # set bullet points
            # slide_subtitle.text = slide_chunk["subtitle"]
            # text_box = slide_slide.shapes.add_textbox(left, top, width, height)
            # for j in range(len(slide_chunk["bullets"])):
            #     tb = text_box.text_frame
            #     prg = tb.add_paragraph()
            #     prg.text = slide_chunk["bullets"][j]

        return self.prs